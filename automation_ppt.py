from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE


prs                = Presentation()
#---------------SLIDE_LAYOUT  N°0 Hello World-------------------------
title_slide_layout = prs.slide_layouts[0]
#---------------PRS-ADD ->TITLE-SLIDE_LAYOUT--------------
slide0              = prs.slides.add_slide(title_slide_layout) # obj slide para agregar 

title              = slide0.shapes.title    # obj title para escribir
subtitle           = slide0.placeholders[1] # obj subtitle para escribir
#---------------text--------------------------------------
title.text         = "Hello, World!"
subtitle.text      = "python-pptx was here, hello world"

#---------------SLIDE_LAYOUT  N°1 EJEMPLO DE DIAPOSITIVA-------------------------
bullet_slide_layout = prs.slide_layouts[1]
#---------------PRS-ADD ->bullet-SLIDE_LAYOUT--------------
slide1      = prs.slides.add_slide(bullet_slide_layout)
title1      = slide1.shapes.title
body_shape1 = slide1.shapes.placeholders[1]
#--------------text----------------------------------------
title1.text= "Adding a Bullet Slide" 
tf         = body_shape1.text_frame

tf.text    = "Find the bullet slide layout"

p      = tf.add_paragraph()
p.text = "Use _TextFrame.text for first bullet"
p.level= 1

p      = tf.add_paragraph()
p.text = "Use _TextFrame.add_paragraoh() for  subsequent bullets"
p.level= 2

#---------------SLIDE_LAYOUT  N°2 IMAGENES-------------------------
blank_slide_layout = prs.slide_layouts[2]
#---------------PRS-ADD ->BLANK-SLIDE_LAYOUT--------------
slide2 = prs.slides.add_slide(blank_slide_layout)
img_path = 'amisr3.png'

# IMAGEN TAMANO NORMAL SIN RESTRICCION
left = top = Inches(1) #position
pic = slide2.shapes.add_picture(img_path, left, top)
# IMAGEN PEQUEÑA
left = Inches(5)  # position
height = Inches(5.5) # tamano
pic = slide2.shapes.add_picture(img_path, left, top, height=height)


#---------------SLIDE_LAYOUT  N°3 FORMAS-------------------------
title_only_slide_layout = prs.slide_layouts[3]
#---------------PRS-ADD ->TITLE-ONLY-SLIDE_LAYOUT--------------

slide3 = prs.slides.add_slide(title_only_slide_layout)
title3 = slide3.shapes.title

#------------------------text------------------------------------
title3.text = "Adding an AutoShape"
#----------------------..forma-----------------------------------
#---position---
left = Inches(0.93)# position x
top  = Inches(3) # position y
#---tamano----
width= Inches(1.75) 
height=Inches(1.0)
shape = slide3.shapes.add_shape(MSO_SHAPE.PENTAGON,left,top,width,height)
shape.text ='Step 1'

left = left + width - Inches(0.4) # position 0.9+1.75-0.4
width = Inches(2.0)  # chevrons need more width for visual balance
for n in range(2, 6):
    shape = slide3.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = 'Step %d' % n
    left = left + width - Inches(0.4)



#---------------SLIDE_LAYOUT  N°4 TABLAS-------------------------
title_only_table_slide_layout = prs.slide_layouts[4]
#---------------PRS-ADD ->TITLE-ONLY-SLIDE_LAYOUT--------------
slide4 = prs.slides.add_slide(title_only_table_slide_layout)
title4 = slide4.shapes.title
#------------------------text------------------------------------
title4.text = "Adding a Table"
#------------------------table-----------------------------------
rows= cols = 2
left= top  = Inches(2) #position
width  = Inches(6.0) # tamano
height = Inches(0.8) # tamano
#---------------------objeto-tabla--------
table = slide4.shapes.add_table(rows,cols,left,top,width,height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Columna 0'
table.cell(0, 1).text = 'Columna 1'

# write body cells
table.cell(1, 0).text = 'Valor 0'
table.cell(1, 1).text = 'Valor 1'

prs.save('test5.pptx')
